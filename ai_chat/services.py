import hashlib
import io
import json
import logging
import re
import time
import zipfile
from pathlib import Path

import requests
from django.core.files.base import ContentFile

logger = logging.getLogger(__name__)

INSTRUCTIONS_PATH = Path(__file__).resolve().parent / 'instructions' / 'system_instructions.md'
# Admin drops reference files here manually — tab1.md…tab4.md are read
# directly as per-tab instructions; everything else gets extracted,
# chunked, and embedded by `python manage.py index_sources`.
SOURCE_DIR = Path(__file__).resolve().parent / 'source'
TAB_INSTRUCTION_FILENAMES = {
    'build_model': 'build_model.md',
    'run_analysis': 'run_analysis.md',
    'post_processing': 'post_processing.md',
    'design': 'design.md',
}
SUPPORTED_SOURCE_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.dxf', '.json', '.txt', '.md', '.csv', '.py'}
FILES_BLOCK_RE = re.compile(r'<<<FILES>>>(.*?)<<<END_FILES>>>', re.DOTALL)

REQUEST_TIMEOUT = 60

# Static catalog of selectable models per provider, tagged free/paid so the
# console can filter the Model dropdown. Used only as a fallback when a live
# fetch isn't possible (no key yet) or fails — see get_models_for_provider().
MODEL_CATALOG = {
    'anthropic': [
        {'id': 'claude-opus-4-8', 'label': 'claude-opus-4-8', 'tier': 'paid'},
        {'id': 'claude-sonnet-5', 'label': 'claude-sonnet-5', 'tier': 'paid'},
        {'id': 'claude-haiku-4-5', 'label': 'claude-haiku-4-5', 'tier': 'free'},
    ],
    'openai': [
        {'id': 'gpt-5', 'label': 'gpt-5', 'tier': 'paid'},
        {'id': 'gpt-5-mini', 'label': 'gpt-5-mini', 'tier': 'paid'},
        {'id': 'gpt-4o-mini', 'label': 'gpt-4o-mini', 'tier': 'free'},
        {'id': 'gpt-3.5-turbo', 'label': 'gpt-3.5-turbo', 'tier': 'free'},
    ],
    'google': [
        {'id': 'gemini-2.5-pro', 'label': 'gemini-2.5-pro', 'tier': 'paid'},
        {'id': 'gemini-2.5-flash', 'label': 'gemini-2.5-flash', 'tier': 'free'},
        {'id': 'gemini-1.5-flash', 'label': 'gemini-1.5-flash', 'tier': 'free'},
    ],
    'mistral': [
        {'id': 'mistral-large-latest', 'label': 'mistral-large-latest', 'tier': 'paid'},
        {'id': 'mistral-small-latest', 'label': 'mistral-small-latest', 'tier': 'free'},
        {'id': 'open-mistral-7b', 'label': 'open-mistral-7b', 'tier': 'free'},
    ],
}


def get_models_for_provider(provider, tier=None, api_key=None):
    """Return the model list for a provider, plus where it came from.

    Returns (models, meta) where meta = {'source': 'live'|'catalog', 'error': str|None}
    so the caller (and the UI) can tell a real live fetch apart from the
    static fallback instead of the two looking identical.
    """
    if api_key:
        live, error = fetch_live_models(provider, api_key)
        if live:
            return live, {'source': 'live', 'error': None}
        models = MODEL_CATALOG.get(provider, [])
        if tier and tier != 'all':
            models = [m for m in models if m['tier'] == tier]
        return models, {'source': 'catalog', 'error': error or 'Live fetch returned no models.'}

    models = MODEL_CATALOG.get(provider, [])
    if tier and tier != 'all':
        models = [m for m in models if m['tier'] == tier]
    return models, {'source': 'catalog', 'error': None}


def fetch_live_models(provider, api_key):
    """Ask the provider's own API which models this key can use.

    Never raises — any network/auth/parsing problem is caught, logged, and
    returned as (=, error_message) so the caller can fall back to the
    static catalog *and* the frontend can show why the live list didn't
    load, instead of the two cases looking identical.
    """
    try:
        if provider == 'anthropic':
            return _fetch_anthropic_models(api_key), None
        if provider == 'openai':
            return _fetch_openai_models(api_key), None
        if provider == 'google':
            return _fetch_google_models(api_key), None
        if provider == 'mistral':
            return _fetch_mistral_models(api_key), None
        return [], f'Unknown provider "{provider}".'
    except requests.RequestException as exc:
        logger.warning('Live model fetch failed for %s: network error: %s', provider, exc)
        return [], f'Could not reach {provider}: {exc}'
    except Exception as exc:  # noqa: BLE001 - never let a bad response 500 the request
        logger.warning('Live model fetch failed for %s: %s', provider, exc, exc_info=True)
        return [], f'{provider} returned an unexpected response: {exc}'


def _raise_models_error(resp, provider):
    if resp.status_code >= 400:
        raise RuntimeError(f'{provider} error {resp.status_code}: {resp.text[:300]}')


def _fetch_anthropic_models(api_key):
    """GET /v1/models, paginated via after_id/has_more (default page size
    is 20, so without pagination a growing catalog gets silently truncated)."""
    out = []
    after_id = None
    while True:
        params = {'limit': 1000}
        if after_id:
            params['after_id'] = after_id
        resp = requests.get(
            'https://api.anthropic.com/v1/models',
            headers={'x-api-key': api_key, 'anthropic-version': '2023-06-01'},
            params=params,
            timeout=REQUEST_TIMEOUT,
        )
        _raise_models_error(resp, 'Anthropic')
        data = resp.json()
        out.extend(
            {'id': m['id'], 'label': m.get('display_name', m['id']), 'tier': 'live'}
            for m in data.get('data', [])
        )
        if not data.get('has_more') or not data.get('last_id'):
            break
        after_id = data['last_id']
    return out


def _fetch_openai_models(api_key):
    resp = requests.get(
        'https://api.openai.com/v1/models',
        headers={'Authorization': f'Bearer {api_key}'},
        timeout=REQUEST_TIMEOUT,
    )
    _raise_models_error(resp, 'OpenAI')
    data = resp.json()
    ids = sorted(m['id'] for m in data.get('data', []) if m.get('id'))
    return [{'id': i, 'label': i, 'tier': 'live'} for i in ids]


def _fetch_google_models(api_key):
    """GET /v1beta/models, paginated via pageToken/nextPageToken (default
    page size is 50, so without pagination the list gets cut off)."""
    out = []
    page_token = None
    while True:
        params = {'key': api_key, 'pageSize': 1000}
        if page_token:
            params['pageToken'] = page_token
        resp = requests.get(
            'https://generativelanguage.googleapis.com/v1beta/models',
            params=params,
            timeout=REQUEST_TIMEOUT,
        )
        _raise_models_error(resp, 'Google')
        data = resp.json()
        for m in data.get('models', []):
            # names come back as 'models/gemini-2.5-pro'
            model_id = (m.get('name') or '').split('/')[-1]
            if model_id:
                out.append({'id': model_id, 'label': m.get('displayName', model_id), 'tier': 'live'})
        page_token = data.get('nextPageToken')
        if not page_token:
            break
    return out


def _fetch_mistral_models(api_key):
    resp = requests.get(
        'https://api.mistral.ai/v1/models',
        headers={'Authorization': f'Bearer {api_key}'},
        timeout=REQUEST_TIMEOUT,
    )
    _raise_models_error(resp, 'Mistral')
    data = resp.json()
    return [{'id': m['id'], 'label': m['id'], 'tier': 'live'} for m in data.get('data', []) if m.get('id')]


def load_system_instructions():
    try:
        return INSTRUCTIONS_PATH.read_text(encoding='utf-8')
    except FileNotFoundError:
        return 'You are a helpful assistant.'


class AIProviderError(Exception):
    pass


def call_ai(provider, model, api_key, system_prompt, history):
    """history: list of {'role': 'user'|'assistant', 'content': str}.
    Returns the raw assistant text (may still contain a <<<FILES>>> block).
    Never raises — any provider/network/parsing failure becomes a plain
    AIProviderError with a readable message, so a bad key or a malformed
    provider response always turns into a normal chat reply instead of a
    server crash.
    """
    try:
        if provider == 'anthropic':
            return _call_anthropic(model, api_key, system_prompt, history)
        if provider == 'openai':
            return _call_openai(model, api_key, system_prompt, history)
        if provider == 'google':
            return _call_google(model, api_key, system_prompt, history)
        if provider == 'mistral':
            return _call_mistral(model, api_key, system_prompt, history)
        raise AIProviderError(f'Unknown provider "{provider}".')
    except AIProviderError:
        raise
    except requests.RequestException as exc:
        raise AIProviderError(f'Could not reach {provider}: {exc}') from exc
    except (KeyError, IndexError, TypeError, ValueError) as exc:
        raise AIProviderError(f'{provider} returned an unexpected response: {exc}') from exc


def _raise_for_status(resp, provider):
    if resp.status_code >= 400:
        raise AIProviderError(f'{provider} error {resp.status_code}: {resp.text[:400]}')


# Transient "the provider is momentarily overloaded" statuses — worth an
# automatic retry rather than immediately failing the chat message.
# 429 = rate limited, 500/502/503 = generic server-side overload,
# 529 = Anthropic's own "overloaded" code.
_RETRYABLE_STATUSES = {429, 500, 502, 503, 529}
_MAX_RETRIES = 2
_RETRY_BACKOFF_SECONDS = 1.5


def _post_with_retry(url, provider, **kwargs):
    """POST with automatic retry (exponential backoff) on transient
    overload/rate-limit responses. Returns the final response either way —
    caller still runs it through _raise_for_status."""
    attempt = 0
    while True:
        resp = requests.post(url, timeout=REQUEST_TIMEOUT, **kwargs)
        if resp.status_code not in _RETRYABLE_STATUSES or attempt >= _MAX_RETRIES:
            return resp
        wait = _RETRY_BACKOFF_SECONDS * (2 ** attempt)
        logger.warning(
            '%s returned %s (attempt %d/%d) — retrying in %.1fs',
            provider, resp.status_code, attempt + 1, _MAX_RETRIES, wait,
        )
        time.sleep(wait)
        attempt += 1


def _call_anthropic(model, api_key, system_prompt, history):
    resp = _post_with_retry(
        'https://api.anthropic.com/v1/messages',
        'Anthropic',
        headers={
            'x-api-key': api_key,
            'anthropic-version': '2023-06-01',
            'content-type': 'application/json',
        },
        json={
            'model': model,
            'max_tokens': 4096,
            'system': system_prompt,
            'messages': [{'role': m['role'], 'content': m['content']} for m in history],
        },
    )
    _raise_for_status(resp, 'Anthropic')
    data = resp.json()
    return ''.join(block.get('text', '') for block in data.get('content', []) if block.get('type') == 'text')


def _call_openai(model, api_key, system_prompt, history):
    resp = _post_with_retry(
        'https://api.openai.com/v1/chat/completions',
        'OpenAI',
        headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
        json={
            'model': model,
            'messages': [{'role': 'system', 'content': system_prompt}] + history,
        },
    )
    _raise_for_status(resp, 'OpenAI')
    data = resp.json()
    return data['choices'][0]['message']['content']


def _call_mistral(model, api_key, system_prompt, history):
    resp = _post_with_retry(
        'https://api.mistral.ai/v1/chat/completions',
        'Mistral',
        headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
        json={
            'model': model,
            'messages': [{'role': 'system', 'content': system_prompt}] + history,
        },
    )
    _raise_for_status(resp, 'Mistral')
    data = resp.json()
    return data['choices'][0]['message']['content']


def _call_google(model, api_key, system_prompt, history):
    contents = [
        {'role': 'model' if m['role'] == 'assistant' else 'user', 'parts': [{'text': m['content']}]}
        for m in history
    ]
    resp = _post_with_retry(
        f'https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent',
        'Google',
        params={'key': api_key},
        json={
            'system_instruction': {'parts': [{'text': system_prompt}]},
            'contents': contents,
        },
    )
    _raise_for_status(resp, 'Google')
    data = resp.json()
    candidates = data.get('candidates') or []
    if not candidates:
        reason = data.get('promptFeedback', {}).get('blockReason', 'no candidates returned')
        raise AIProviderError(f'Google returned no response ({reason}).')
    parts = candidates[0].get('content', {}).get('parts', [])
    return ''.join(p.get('text', '') for p in parts)


def extract_file_directives(raw_text):
    """Split the model's raw reply into (visible_reply, [{'filename','content'}...])."""
    match = FILES_BLOCK_RE.search(raw_text)
    if not match:
        return raw_text.strip(), []

    visible = (raw_text[:match.start()] + raw_text[match.end():]).strip()
    try:
        files = json.loads(match.group(1).strip())
        if not isinstance(files, list):
            files = []
    except (json.JSONDecodeError, ValueError):
        files = []

    cleaned = []
    for f in files:
        name = str(f.get('filename', '')).strip().replace('/', '_').replace('\\', '_')
        content = f.get('content', '')
        if name:
            cleaned.append({'filename': name, 'content': str(content)})
    return visible, cleaned


def _build_txt_like(content):
    return content.encode('utf-8')


def _build_docx(content):
    from docx import Document
    doc = Document()
    for line in content.splitlines() or ['']:
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx(content):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for line in content.splitlines():
        delim = '\t' if '\t' in line else ','
        ws.append(line.split(delim))
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx(content):
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]
    lines = content.splitlines() or ['']
    for line in lines:
        title, _, body = line.partition('|')
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title.strip() or 'Slide'
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = body.strip()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_dxf(content):
    import ezdxf
    doc = ezdxf.new()
    msp = doc.modelspace()
    y = 0
    for line in content.splitlines() or ['']:
        msp.add_text(line, dxfattribs={'height': 2.5}).set_placement((0, y))
        y -= 5
    buf = io.StringIO()
    doc.write(buf)
    return buf.getvalue().encode('utf-8')


_BUILDERS = {
    '.txt': _build_txt_like,
    '.py': _build_txt_like,
    '.json': _build_txt_like,
    '.md': _build_txt_like,
    '.csv': _build_txt_like,
    '.docx': _build_docx,
    '.xlsx': _build_xlsx,
    '.pptx': _build_pptx,
    '.dxf': _build_dxf,
}


def build_file_bytes(filename, content):
    ext = Path(filename).suffix.lower()
    builder = _BUILDERS.get(ext, _build_txt_like)
    try:
        return builder(content)
    except Exception:
        # Missing optional dependency or malformed content — fall back to
        # a plain-text artifact so the request never hard-fails.
        return content.encode('utf-8')


def zip_files(files):
    """files: [{'filename','content'}...] -> (zip_bytes, [names])."""
    buf = io.BytesIO()
    names = []
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        used = set()
        for f in files:
            name = f['filename']
            base, count = name, 1
            while name in used:
                stem = Path(base).stem
                suffix = Path(base).suffix
                name = f'{stem}_{count}{suffix}'
                count += 1
            used.add(name)
            zf.writestr(name, build_file_bytes(name, f['content']))
            names.append(name)
    return buf.getvalue(), names


def save_output_bundle(project, tab, message, files, zip_basename):
    zip_bytes, names = zip_files(files)
    from .models import OutputBundle
    bundle = OutputBundle(project=project, tab=tab, message=message, file_names=names, size_bytes=len(zip_bytes))
    bundle.zip_file.save(zip_basename, ContentFile(zip_bytes), save=True)
    return bundle


# ============================================================
# Source folder (admin "training" material) — extraction, chunking,
# embedding, and retrieval. Admin manually places files in SOURCE_DIR:
#   ai_chat/source/tab1.md ... tab4.md   -> per-tab instructions
#   ai_chat/source/<anything else>       -> indexed as retrievable material
# Run `python manage.py index_sources` after adding/changing/removing
# files. See ai_chat/models.py: SourceDocument, SourceChunk,
# EmbeddingSettings.
# ============================================================

def extract_text_from_source(data, filename):
    """data: raw file bytes. Returns extracted plain text, or raises
    ValueError for an unsupported extension."""
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    if ext == 'pdf':
        return _extract_pdf(data)
    if ext == 'docx':
        return _extract_docx(data)
    if ext in ('xlsx', 'xls'):
        return _extract_xlsx(data)
    if ext == 'pptx':
        return _extract_pptx(data)
    if ext == 'dxf':
        return _extract_dxf(data)
    if ext == 'json':
        return _extract_json(data)
    if ext in ('txt', 'md', 'csv', 'py'):
        return data.decode('utf-8', errors='replace')
    raise ValueError(f'Unsupported source file type: .{ext or "?"}')


def _extract_pdf(data):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return '\n\n'.join((page.extract_text() or '') for page in reader.pages)


def _extract_docx(data):
    import docx
    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(' | '.join(cell.text for cell in row.cells))
    return '\n'.join(parts)


def _extract_xlsx(data):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f'## Sheet: {sheet.title}')
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(' | '.join(cells))
    return '\n'.join(lines)


def _extract_pptx(data):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'## Slide {i}')
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
    return '\n'.join(lines)


def _extract_dxf(data):
    import ezdxf
    doc = ezdxf.read(io.StringIO(data.decode('utf-8', errors='replace')))
    msp = doc.modelspace()
    lines = []
    for entity in msp:
        if entity.dxftype() == 'TEXT':
            lines.append(entity.dxf.text)
        elif entity.dxftype() == 'MTEXT':
            lines.append(entity.plain_text() if hasattr(entity, 'plain_text') else entity.text)
    return '\n'.join(lines)


def _extract_json(data):
    obj = json.loads(data.decode('utf-8', errors='replace'))
    return json.dumps(obj, indent=2, ensure_ascii=False)


def chunk_text(text, chunk_size=1200, overlap=200):
    """Simple character-based chunker with overlap. Good enough for
    admin-provided reference material at moderate scale; if documents get
    very large or retrieval quality needs improving, a sentence/paragraph-
    aware chunker would be the next upgrade."""
    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - overlap
    return chunks


def embed_texts(texts, batch_size=64):
    """Embed a list of strings with the admin-configured embedding
    provider/model (see EmbeddingSettings, set in /admin/). Returns one
    vector per input text, same order. Batches requests so a large
    document doesn't blow past a provider's per-request input limit."""
    from .models import EmbeddingSettings
    row = EmbeddingSettings.load()
    api_key = row.get_key()
    if not api_key:
        raise AIProviderError(
            'No embedding API key configured — set one in /admin/ under Embedding settings.'
        )

    embedder = {
        'openai': _embed_openai,
        'google': _embed_google,
        'mistral': _embed_mistral,
    }.get(row.provider)
    if embedder is None:
        raise AIProviderError(f'Unsupported embedding provider: {row.provider}')

    out = []
    for i in range(0, len(texts), batch_size):
        out.extend(embedder(texts[i:i + batch_size], api_key, row.model))
    return out


def _embed_openai(texts, api_key, model):
    resp = _post_with_retry(
        'https://api.openai.com/v1/embeddings',
        'OpenAI',
        headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
        json={'model': model, 'input': texts},
    )
    _raise_for_status(resp, 'OpenAI')
    data = resp.json()
    return [item['embedding'] for item in sorted(data['data'], key=lambda d: d['index'])]


def _embed_google(texts, api_key, model):
    resp = _post_with_retry(
        f'https://generativelanguage.googleapis.com/v1beta/models/{model}:batchEmbedContents',
        'Google',
        params={'key': api_key},
        json={
            'requests': [
                {'model': f'models/{model}', 'content': {'parts': [{'text': t}]}}
                for t in texts
            ]
        },
    )
    _raise_for_status(resp, 'Google')
    data = resp.json()
    return [e['values'] for e in data.get('embeddings', [])]


def _embed_mistral(texts, api_key, model):
    resp = _post_with_retry(
        'https://api.mistral.ai/v1/embeddings',
        'Mistral',
        headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
        json={'model': model, 'input': texts},
    )
    _raise_for_status(resp, 'Mistral')
    data = resp.json()
    return [item['embedding'] for item in sorted(data['data'], key=lambda d: d['index'])]


def _file_hash(data):
    return hashlib.sha256(data).hexdigest()


def list_source_files():
    """Every file directly inside SOURCE_DIR that isn't a tab*.md
    instruction file and has a supported extension. Returns relative
    filenames (str)."""
    if not SOURCE_DIR.is_dir():
        return []
    names = []
    for path in sorted(SOURCE_DIR.iterdir()):
        if not path.is_file():
            continue
        if path.name in TAB_INSTRUCTION_FILENAMES.values():
            continue
        if path.suffix.lower() not in SUPPORTED_SOURCE_EXTENSIONS:
            continue
        names.append(path.name)
    return names


def load_tab_instructions(tab_key):
    """Read ai_chat/source/tab<N>.md for this tab, or '' if it doesn't
    exist yet — a missing file just means this tab has no extra
    instructions beyond the base system prompt."""
    filename = TAB_INSTRUCTION_FILENAMES.get(tab_key)
    if not filename:
        return ''
    path = SOURCE_DIR / filename
    if not path.is_file():
        return ''
    return path.read_text(encoding='utf-8').strip()


def index_sources():
    """Scan SOURCE_DIR and bring SourceDocument/SourceChunk in sync with
    whatever's actually on disk:
      - new or changed (by content hash) files are extracted, chunked,
        embedded, and (re)stored
      - unchanged files are skipped
      - files that no longer exist in the folder are removed from the DB
    Returns a dict summarizing what happened, for a management command
    (or an admin action) to report back."""
    from django.utils import timezone
    from .models import SourceChunk, SourceDocument

    result = {'added': [], 'updated': [], 'skipped': [], 'removed': [], 'errors': []}

    on_disk = set(list_source_files())
    existing = {doc.relative_path: doc for doc in SourceDocument.objects.all()}

    # Files that disappeared from the folder since the last scan.
    for relative_path, doc in existing.items():
        if relative_path not in on_disk:
            doc.delete()
            result['removed'].append(relative_path)

    for relative_path in sorted(on_disk):
        path = SOURCE_DIR / relative_path
        data = path.read_bytes()
        file_hash = _file_hash(data)
        doc = existing.get(relative_path)
        is_new = doc is None
        if doc is None:
            doc = SourceDocument(relative_path=relative_path)
        elif doc.status == 'processed' and doc.file_hash == file_hash:
            result['skipped'].append(relative_path)
            continue

        try:
            text = extract_text_from_source(data, relative_path)
            chunks = chunk_text(text)
            if not chunks:
                raise ValueError('No extractable text found in this file.')
            embeddings = embed_texts(chunks)
        except Exception as exc:
            doc.file_hash = file_hash
            doc.status = 'error'
            doc.processing_error = str(exc)
            doc.chunk_count = 0
            doc.save()
            result['errors'].append((relative_path, str(exc)))
            continue

        doc.file_hash = file_hash
        doc.status = 'processed'
        doc.processing_error = ''
        doc.chunk_count = len(chunks)
        doc.indexed_at = timezone.now()
        doc.save()

        doc.chunks.all().delete()
        SourceChunk.objects.bulk_create([
            SourceChunk(document=doc, chunk_index=i, content=chunk, embedding=embedding)
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
        ])

        (result['added'] if is_new else result['updated']).append(relative_path)

    return result


def _cosine_similarity(a, b):
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = sum(x * x for x in a) ** 0.5
    norm_b = sum(y * y for y in b) ** 0.5
    if not norm_a or not norm_b:
        return 0.0
    return dot / (norm_a * norm_b)


def retrieve_relevant_chunks(query, top_k=5):
    """Embed the query and return the top_k most similar SourceChunks
    (across every indexed file — sources are shared across tabs). Never
    raises: an embedding failure just means no grounding material for
    this reply, so the chat keeps working."""
    from .models import SourceChunk
    chunks = list(SourceChunk.objects.select_related('document'))
    if not chunks:
        return []

    try:
        query_embedding = embed_texts([query])[0]
    except Exception as exc:
        logger.warning('Query embedding failed, skipping retrieval: %s', exc, exc_info=True)
        return []

    scored = [(_cosine_similarity(query_embedding, c.embedding), c) for c in chunks]
    scored.sort(key=lambda pair: pair[0], reverse=True)
    return [c for score, c in scored[:top_k] if score > 0]


def build_system_prompt(tab_key, user_query):
    """Base instructions + this tab's tab<N>.md instructions + the most
    relevant chunks of admin-provided source material for user_query."""
    parts = [load_system_instructions()]

    tab_text = load_tab_instructions(tab_key)
    if tab_text:
        parts.append(f'## Instructions for this tab\n{tab_text}')

    relevant = retrieve_relevant_chunks(user_query) if user_query else []
    if relevant:
        material = '\n\n'.join(
            f'[Source: {c.document.relative_path}]\n{c.content}' for c in relevant
        )
        parts.append(
            '## Reference material\n'
            "Use these excerpts from admin-provided source documents to ground your "
            "answer where relevant. If they don't contain the answer, say so rather "
            "than guessing.\n\n" + material
        )

    return '\n\n---\n\n'.join(parts)
